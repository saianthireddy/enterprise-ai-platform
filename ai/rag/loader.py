"""Document loaders: PDF, DOCX, XLSX, PPTX, and plain text -> raw text."""
from __future__ import annotations

from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def load_pdf(path: str | Path) -> str:
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def load_docx(path: str | Path) -> str:
    doc = DocxDocument(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_xlsx(path: str | Path) -> str:
    wb = load_workbook(str(path), data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def load_pptx(path: str | Path) -> str:
    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def load_text(path: str | Path) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
    ".txt": load_text,
    ".md": load_text,
}


def load_document(path: str | Path) -> str:
    suffix = Path(path).suffix.lower()
    loader = LOADERS.get(suffix)
    if loader is None:
        raise ValueError(f"Unsupported document type: {suffix}")
    return loader(path)
